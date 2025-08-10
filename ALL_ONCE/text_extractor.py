# text_extractor.py
# Description: Functions for extracting text content from various file types.

import io
from typing import List, Dict, Union
from zipfile import ZipFile

import fitz  # PyMuPDF
import docx
from pptx import Presentation
import openpyxl
from PIL import Image
import pytesseract
from bs4 import BeautifulSoup

from file_downloader import run_download_and_identify_file

def _extract_from_html(file_content: bytes) -> List[Dict[str, Union[str, int]]]:
    """Extracts visible text from HTML content."""
    text = ""
    try:
        soup = BeautifulSoup(file_content, 'html.parser')
        for script_or_style in soup(["script", "style"]):
            script_or_style.decompose()
        text = soup.get_text()
        lines = (line.strip() for line in text.splitlines())
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        text = '\n'.join(chunk for chunk in chunks if chunk)
    except Exception as e:
        print(f"Error extracting from HTML: {e}")
    return [{"page": 1, "text": text}]

def _extract_from_pdf(file_content: bytes) -> List[Dict[str, Union[str, int]]]:
    """Extracts text from each page of a PDF."""
    pages_text = []
    try:
        with fitz.open(stream=file_content, filetype="pdf") as doc:
            for i, page in enumerate(doc):
                pages_text.append({"page": i + 1, "text": page.get_text("text", sort=True) or ""})
    except Exception as e:
        print(f"Error extracting from PDF: {e}")
    return pages_text

def _extract_from_docx(file_content: bytes) -> List[Dict[str, Union[str, int]]]:
    """Extracts text from a DOCX document."""
    text = ""
    try:
        doc = docx.Document(io.BytesIO(file_content))
        for para in doc.paragraphs:
            text += para.text + "\n"
    except Exception as e:
        print(f"Error extracting from DOCX: {e}")
    return [{"page": 1, "text": text}]

def _extract_from_pptx(file_content: bytes) -> List[Dict[str, Union[str, int]]]:
    """Extracts text and performs OCR on images from a PPTX presentation."""
    pages_text = []
    try:
        prs = Presentation(io.BytesIO(file_content))
        for i, slide in enumerate(prs.slides):
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + "\n"
                if hasattr(shape, "image"):
                    try:
                        image_bytes = shape.image.blob
                        image = Image.open(io.BytesIO(image_bytes))
                        ocr_text = pytesseract.image_to_string(image)
                        if ocr_text.strip():
                            slide_text += f"\n--- OCR Text from Image ---\n{ocr_text.strip()}\n--- End OCR Text ---\n"
                    except Exception as e:
                        print(f"Warning: Could not process an image on slide {i + 1}. Error: {e}")
            if slide_text.strip():
                pages_text.append({"page": i + 1, "text": slide_text})
    except Exception as e:
        print(f"Error extracting from PPTX: {e}")
    return pages_text

def _extract_from_xlsx(file_content: bytes) -> List[Dict[str, Union[str, int]]]:
    """Extracts text from each sheet of an XLSX workbook."""
    pages_text = []
    try:
        workbook = openpyxl.load_workbook(io.BytesIO(file_content))
        for i, sheet_name in enumerate(workbook.sheetnames):
            sheet_text = ""
            sheet = workbook[sheet_name]
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value:
                        sheet_text += str(cell.value) + " "
                sheet_text += "\n"
            pages_text.append({"page": i + 1, "text": f"Sheet: {sheet_name}\n{sheet_text}"})
    except Exception as e:
        print(f"Error extracting from XLSX: {e}")
    return pages_text

def _extract_from_image(file_content: bytes) -> List[Dict[str, Union[str, int]]]:
    """Extracts text from an image file using OCR."""
    text = ""
    try:
        image = Image.open(io.BytesIO(file_content))
        text = pytesseract.image_to_string(image)
    except Exception as e:
        print(f"Error extracting from image with OCR: {e}")
    return [{"page": 1, "text": text}]

def _extract_from_zip(file_content: bytes) -> List[Dict[str, Union[str, int]]]:
    """Recursively extracts text from files within a ZIP archive."""
    all_texts = []
    CONSECUTIVE_FAILURE_THRESHOLD = 4
    consecutive_failures = 0
    try:
        with ZipFile(io.BytesIO(file_content)) as zf:
            for file_info in zf.infolist():
                if file_info.is_dir():
                    continue
                file_name = file_info.filename
                if '.' not in file_name:
                    consecutive_failures += 1
                    print(f"-> Skipping file with no extension: '{file_name}'")
                    continue
                
                file_ext = file_name.split('.')[-1].lower()
                with zf.open(file_info) as file:
                    content = file.read()
                    extracted_texts = extract_text(content, file_ext)

                    if extracted_texts:
                        consecutive_failures = 0
                        for item in extracted_texts:
                            item["text"] = f"From ZIP archive '{file_name}':\n{item['text']}"
                        all_texts.extend(extracted_texts)
                    else:
                        consecutive_failures += 1
                
                if consecutive_failures >= CONSECUTIVE_FAILURE_THRESHOLD:
                    print(f"\n⚠️ ABORTING ZIP EXTRACTION: Detected {consecutive_failures} consecutive unsupported files.")
                    return []
    except Exception as e:
        print(f"Error extracting from ZIP: {e}")
        return []
    return all_texts

# A dictionary to map file types to their extraction functions
EXTRACTION_MAP = {
    'html': _extract_from_html,
    'pdf': _extract_from_pdf,
    'docx': _extract_from_docx,
    'pptx': _extract_from_pptx,
    'xlsx': _extract_from_xlsx,
    'png': _extract_from_image,
    'jpg': _extract_from_image,
    'jpeg': _extract_from_image,
    'zip': _extract_from_zip,
}

def extract_text(file_content: bytes, file_type: str) -> List[Dict[str, Union[str, int]]]:
    """
    Dispatcher function to extract text based on file type.
    """
    extractor_func = EXTRACTION_MAP.get(file_type)
    if extractor_func:
        return extractor_func(file_content)
    else:
        # Fallback for unknown text-based types
        try:
            text = file_content.decode('utf-8')
            return [{"page": 1, "text": text}]
        except UnicodeDecodeError:
            print(f"Unsupported and non-text file type: {file_type}")
            return []
        
# if __name__ == "__main___":
    # pdf_url = "https://hackrx.blob.core.windows.net/assets/policy.pdf?sv=2023-01-03&st=2025-07-04T09%3A11%3A24Z&se=2027-07-05T09%3A11%3A00Z&sr=b&sp=r&sig=N4a9OU0w0QXO6AOIBiu4bpl7AXvEZogeT%2FjUHNO7HzQ%3D"
    # file_content, file_type = run_download_and_identify_file(pdf_url)

    # content = extract_text(file_content, file_type)
    # print(content[:5])